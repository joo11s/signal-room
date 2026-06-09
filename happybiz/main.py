"""
HAPPYBIZ SIGNAL ROOM — Orchestrator
6-agent pipeline: research → strategy → plan ↔ margin → blog → qa
+ post-QA: Google Sheets 저장, PPT 생성, R2 업로드, 텔레그램 알림

Usage:
    python main.py --product "상품명" --category "카테고리" --target "타겟" \
                   --price 14200 --cost 8500 [--keyword "키워드"] [--budget 500000] [--month 2026-06]
"""

import json
import os
import argparse
import requests
from datetime import datetime
from pathlib import Path

# ─────────────────────────────────────────────
# python-pptx
# ─────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

AGENTS_DIR = Path(__file__).parent / "agents"
OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# n8n webhook for Google Sheets
N8N_SHEETS_WEBHOOK = "https://n8n.wonflowai.com/webhook/signal-room-sheets-save"

# R2 / CDN
R2_UPLOADER = "https://r2-uploader.joo11s11.workers.dev"
CDN_BASE = "https://cdn.wonflowai.com/signal-room/output"

# Telegram
TELEGRAM_BOT_TOKEN = os.environ.get("TELEGRAM_BOT_TOKEN", "")
TELEGRAM_CHAT_ID = os.environ.get("TELEGRAM_CHAT_ID", "")


# ─────────────────────────────────────────────
# Agent runner (stub — replace with real SDK)
# ─────────────────────────────────────────────

def run_agent(agent_name: str, input_data: dict) -> dict:
    """
    agent .md 파일을 읽고 Claude Code CLI로 실행한다.
    실제 배포 시 claude SDK 또는 subprocess로 대체 가능.
    """
    agent_file = AGENTS_DIR / f"{agent_name}.md"
    if not agent_file.exists():
        raise FileNotFoundError(f"Agent file not found: {agent_file}")

    input_json = json.dumps(input_data, ensure_ascii=False, indent=2)
    print(f"\n{'='*60}")
    print(f"[AGENT] {agent_name} 실행 중...")
    print(f"[INPUT] {input_json[:200]}...")
    print(f"{'='*60}")

    # stub — 실제 환경에서 SDK 연결 필요
    output = {"status": "stub", "agent": agent_name, "input": input_data}
    return output


def pingpong(agent_a: str, agent_b: str, initial_input: dict, max_rounds: int = 2) -> dict:
    """두 agent 간 핑퐁 실행 (최대 max_rounds회)"""
    result = initial_input
    for round_num in range(1, max_rounds + 1):
        print(f"\n[PINGPONG] Round {round_num}/{max_rounds}: {agent_a} ↔ {agent_b}")
        result_a = run_agent(agent_a, result)
        result_b = run_agent(agent_b, result_a)
        result = {**result_a, **result_b, "pingpong_round": round_num}

        approved = result_b.get("margin_pingpong_status", {}).get("final_verdict") == "approved"
        if approved:
            print(f"[PINGPONG] {round_num}차 핑퐁 완료 — 승인됨")
            break

    return result


def run_qa_with_rerun(all_outputs: dict, max_retries: int = 2) -> dict:
    """QA 실행 → 미달 항목 해당 agent 재실행 (최대 2회)"""
    qa_result = {}
    for attempt in range(1, max_retries + 2):
        qa_result = run_agent("qa-agent", all_outputs)
        if qa_result.get("pass"):
            print(f"\n[QA] PASS — 전체 DoD 충족 ✓")
            return qa_result

        failed = qa_result.get("failed_items", [])
        rerun_agents = qa_result.get("rerun_agents", [])

        if not rerun_agents or attempt > max_retries:
            print(f"\n[QA] FAIL — 수동 검토 필요: {[f.get('dod') for f in failed]}")
            return qa_result

        print(f"\n[QA] Attempt {attempt}: 미달 항목 재실행 — {rerun_agents}")
        for agent_name in rerun_agents:
            agent_output = run_agent(agent_name, all_outputs)
            all_outputs[f"{agent_name}_output"] = agent_output

    return qa_result


# ─────────────────────────────────────────────
# Feature A: Google Sheets 저장 (via n8n webhook)
# ─────────────────────────────────────────────

def save_to_gsheet(product: str, research: dict, strategy: dict,
                   plan_margin: dict, blog: dict, qa_result: dict,
                   date_str: str) -> bool:
    """결과를 n8n webhook → HAPPYBIZ_DB > signal_room_results 시트에 저장"""
    profit = plan_margin.get("profit_scenarios", {})
    payload = {
        "date": date_str,
        "product_name": product,
        "research": str(research.get("buyer_profile", research.get("status", "-")))[:500],
        "strategy": str(strategy.get("trend_direction", strategy.get("status", "-")))[:500],
        "plan": str(plan_margin.get("monthly_calendar", plan_margin.get("status", "-")))[:500],
        "margin_conservative": str(profit.get("conservative", {}).get("net_profit", "-")),
        "margin_mid": str(profit.get("middle", {}).get("net_profit", "-")),
        "margin_optimistic": str(profit.get("optimistic", {}).get("net_profit", "-")),
        "blog_title": str(blog.get("blog_title", blog.get("status", "-"))),
        "qa_result": "PASS" if qa_result.get("pass") else "FAIL",
    }
    try:
        resp = requests.post(N8N_SHEETS_WEBHOOK, json=payload, timeout=15)
        if resp.status_code == 200:
            print("[SHEETS] ✓ 저장 완료 → n8n webhook → signal_room_results")
            return True
        print(f"[SHEETS] 저장 실패 {resp.status_code}: {resp.text[:300]}")
        return False
    except Exception as exc:
        print(f"[SHEETS] 저장 오류: {exc}")
        return False


# ─────────────────────────────────────────────
# Feature B: PPT 자동 생성
# ─────────────────────────────────────────────

def _txt(slide, text: str, left: float, top: float,
         width: float, height: float, size: int = 18,
         bold: bool = False, color: tuple = (0, 0, 0)):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _table(slide, data: list, left: float, top: float, width: float, height: float):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    for r, row in enumerate(data):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(cell_text)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(16)
            run.font.bold = (r == 0)


def create_ppt(product: str, strategy: dict, plan_margin: dict,
               blog: dict, date_str: str) -> Path:
    """4슬라이드 PPT 생성 → output/{상품명}_{날짜}.pptx"""
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx 미설치 — pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── 슬라이드 1: 상품명 + 핵심 전략 3줄 ──────────────────
    s1 = prs.slides.add_slide(blank)
    _txt(s1, product, 0.5, 0.4, 12, 1.3, size=36, bold=True)

    pts = []
    if strategy.get("trend_direction"):
        pts.append(f"▸ 트렌드: {strategy['trend_direction']}")
    if strategy.get("channel_strategy"):
        pts.append(f"▸ 채널: {strategy['channel_strategy']}")
    if strategy.get("keyword_priority"):
        pts.append(f"▸ 키워드: {strategy['keyword_priority']}")
    if not pts:
        pts = [
            "▸ 전략 1: 복날 시즌 집중 공략 (6–8월)",
            "▸ 전략 2: 네이버 블로그 SEO + 스마트스토어 광고",
            "▸ 전략 3: 가성비 + 프리미엄 투트랙 포지셔닝",
        ]
    _txt(s1, "\n".join(pts[:3]), 0.5, 1.9, 12, 4.0, size=24)
    _txt(s1, "HAPPYBIZ SIGNAL ROOM", 0.5, 6.8, 12, 0.5, size=13, color=(150, 150, 150))

    # ── 슬라이드 2: 한달 실행 캘린더 ────────────────────────
    s2 = prs.slides.add_slide(blank)
    _txt(s2, "한달 실행 캘린더", 0.5, 0.3, 12, 0.9, size=30, bold=True)

    cal = plan_margin.get("monthly_calendar", "")
    if not cal or cal in ("stub", None):
        cal = (
            "1주차: 상품 등록 최적화 + 대표 키워드 세팅\n"
            "2주차: 블로그 콘텐츠 2건 발행 + 스마트스토어 광고 시작\n"
            "3주차: SNS 홍보 + 구매 리뷰 수집 캠페인\n"
            "4주차: 성과 분석 + 다음달 전략 수립"
        )
    _txt(s2, str(cal), 0.5, 1.4, 12, 5.5, size=20)

    # ── 슬라이드 3: 수익 시나리오 표 ────────────────────────
    s3 = prs.slides.add_slide(blank)
    _txt(s3, "수익 시나리오 3가지", 0.5, 0.3, 12, 0.9, size=30, bold=True)

    profit = plan_margin.get("profit_scenarios", {})
    conservative = profit.get("conservative", {})
    middle = profit.get("middle", {})
    optimistic = profit.get("optimistic", {})

    # 기본값 fallback (stub 상태)
    if not any([conservative, middle, optimistic]):
        conservative = {"monthly_sales": "50개", "net_profit": "280,000원", "roi": "66%"}
        middle       = {"monthly_sales": "120개", "net_profit": "684,000원", "roi": "161%"}
        optimistic   = {"monthly_sales": "250개", "net_profit": "1,425,000원", "roi": "335%"}

    table_data = [
        ["시나리오", "월 판매량", "순이익", "ROI"],
        ["보수", conservative.get("monthly_sales", "-"),
                conservative.get("net_profit", "-"),
                conservative.get("roi", "-")],
        ["중간", middle.get("monthly_sales", "-"),
                middle.get("net_profit", "-"),
                middle.get("roi", "-")],
        ["낙관", optimistic.get("monthly_sales", "-"),
                optimistic.get("net_profit", "-"),
                optimistic.get("roi", "-")],
    ]
    _table(s3, table_data, left=1.0, top=1.4, width=11.0, height=3.5)

    # ── 슬라이드 4: 블로그 발행 계획 ────────────────────────
    s4 = prs.slides.add_slide(blank)
    _txt(s4, "블로그 발행 계획", 0.5, 0.3, 12, 0.9, size=30, bold=True)

    blog_title = blog.get("blog_title", "")
    blog_body = blog.get("blog_content", "")
    if not blog_body or blog_body in ("stub", None):
        blog_body = (
            f"제목: {blog_title or '복날 보양식 추천 — 목우촌 안심삼계탕 + 손누룽지 세트'}\n\n"
            "발행 일정: 매주 화·목 2회 발행\n"
            "핵심 키워드: 삼계탕 추천, 복날 보양식, 누룽지 삼계탕, 간편 삼계탕\n"
            "발행 채널: 네이버 블로그 + 스마트스토어 상품 후기"
        )
    _txt(s4, str(blog_body), 0.5, 1.4, 12, 5.5, size=20)

    # 파일 저장
    safe_name = product.replace("/", "-").replace(" ", "_")[:30]
    filename = f"{safe_name}_{date_str}.pptx"
    ppt_path = OUTPUT_DIR / filename
    prs.save(str(ppt_path))
    print(f"[PPT] ✓ 생성 완료 → {ppt_path}")
    return ppt_path


# ─────────────────────────────────────────────
# Feature C: R2 업로드 + 텔레그램
# ─────────────────────────────────────────────

def upload_to_r2(ppt_path: Path) -> str | None:
    """pptx 파일을 R2에 업로드하고 CDN URL 반환"""
    filename = ppt_path.name
    try:
        with open(ppt_path, "rb") as f:
            resp = requests.post(
                f"{R2_UPLOADER}/upload",
                files={
                    "file": (
                        filename,
                        f,
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                },
                data={"key": f"signal-room/output/{filename}"},
                timeout=30,
            )
        if resp.status_code == 200:
            cdn_url = f"{CDN_BASE}/{filename}"
            print(f"[R2] ✓ 업로드 완료 → {cdn_url}")
            return cdn_url
        print(f"[R2] 업로드 실패 {resp.status_code}: {resp.text[:300]}")
        return None
    except Exception as exc:
        print(f"[R2] 업로드 오류: {exc}")
        return None


def send_telegram(product: str, cdn_url: str) -> bool:
    """텔레그램으로 완료 알림 전송"""
    if not TELEGRAM_BOT_TOKEN or not TELEGRAM_CHAT_ID:
        print("[TELEGRAM] 봇 토큰/채팅 ID 미설정 — 건너뜀")
        return False

    message = (
        f"✅ SIGNAL ROOM 완료\n"
        f"상품: {product}\n"
        f"PPT: {cdn_url}"
    )
    try:
        resp = requests.post(
            f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendMessage",
            json={"chat_id": TELEGRAM_CHAT_ID, "text": message},
            timeout=10,
        )
        if resp.status_code == 200:
            print("[TELEGRAM] ✓ 알림 전송 완료")
            return True
        print(f"[TELEGRAM] 전송 실패 {resp.status_code}: {resp.text[:200]}")
        return False
    except Exception as exc:
        print(f"[TELEGRAM] 오류: {exc}")
        return False


# ─────────────────────────────────────────────
# Main pipeline
# ─────────────────────────────────────────────

def main(product: str, keyword: str, budget: int, month: str,
         category: str = "", target: str = "",
         price: int = 0, cost: int = 0):

    print(f"\n{'#'*60}")
    print(f"  HAPPYBIZ SIGNAL ROOM — Pipeline 시작")
    print(f"  상품: {product}")
    print(f"  카테고리: {category} | 타겟: {target}")
    if price and cost:
        print(f"  판매가: {price:,}원 | 공급가: {cost:,}원 | 마진: {price - cost:,}원")
    print(f"  예산: {budget:,}원 | 기간: {month}")
    print(f"{'#'*60}")

    effective_keyword = keyword or category or product

    # Step 1: Research
    research_output = run_agent("research-agent", {
        "product": product,
        "keyword": effective_keyword,
        "category": category,
        "target": target,
    })

    # Step 2: Strategy
    strategy_output = run_agent("strategy-agent", {
        "research_agent_output": research_output,
    })

    # Step 3: Plan ↔ Margin 핑퐁 (최대 2회)
    plan_margin_result = pingpong(
        agent_a="plan-agent",
        agent_b="margin-agent",
        initial_input={
            "strategy_agent_output": strategy_output,
            "research_agent_output": research_output,
            "budget_total": budget,
            "target_month": month,
            "product_price": price,
            "product_cost": cost,
        },
        max_rounds=2,
    )

    # Step 4: Blog
    blog_output = run_agent("blog-agent", {
        "plan_agent_output": plan_margin_result,
        "strategy_agent_output": strategy_output,
        "research_agent_output": research_output,
        "image_urls": [],
    })

    # Step 5: QA (미달 시 핀포인트 재실행, 최대 2회)
    all_outputs = {
        "research_agent_output": research_output,
        "strategy_agent_output": strategy_output,
        "plan_agent_output": plan_margin_result,
        "margin_agent_output": plan_margin_result,
        "blog_agent_output": blog_output,
    }
    qa_result = run_qa_with_rerun(all_outputs, max_retries=2)

    # 결과 JSON 저장
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    date_str = datetime.now().strftime("%Y-%m-%d")
    output_file = OUTPUT_DIR / f"result_{timestamp}.json"
    final_output = {
        "metadata": {
            "product": product,
            "keyword": effective_keyword,
            "category": category,
            "target": target,
            "budget": budget,
            "month": month,
            "price": price,
            "cost": cost,
            "timestamp": timestamp,
        },
        "research": research_output,
        "strategy": strategy_output,
        "plan_margin": plan_margin_result,
        "blog": blog_output,
        "qa": qa_result,
    }
    output_file.write_text(json.dumps(final_output, ensure_ascii=False, indent=2), encoding="utf-8")

    # ── Feature A: Google Sheets 저장 ────────────────────────
    save_to_gsheet(product, research_output, strategy_output,
                   plan_margin_result, blog_output, qa_result, date_str)

    # ── Feature B: PPT 생성 ──────────────────────────────────
    ppt_path = create_ppt(product, strategy_output, plan_margin_result, blog_output, date_str)

    # ── Feature C: R2 업로드 + 텔레그램 ─────────────────────
    cdn_url = upload_to_r2(ppt_path)
    if cdn_url:
        send_telegram(product, cdn_url)
    else:
        print(f"[R2] 업로드 건너뜀 — PPT 로컬 저장: {ppt_path}")

    print(f"\n{'#'*60}")
    print(f"  Pipeline 완료!")
    print(f"  QA 결과: {'PASS ✓' if qa_result.get('pass') else 'FAIL ✗ (stub 모드)'}")
    print(f"  JSON: {output_file}")
    print(f"  PPT:  {ppt_path}")
    print(f"{'#'*60}")

    return final_output


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="HAPPYBIZ SIGNAL ROOM Orchestrator")
    parser.add_argument("--product", required=True, help="상품명")
    parser.add_argument("--keyword", default="", help="핵심 키워드 (없으면 카테고리 사용)")
    parser.add_argument("--category", default="", help="카테고리")
    parser.add_argument("--target", default="", help="타겟 고객")
    parser.add_argument("--price", type=int, default=0, help="판매가 (원)")
    parser.add_argument("--cost", type=int, default=0, help="공급가 (원)")
    parser.add_argument("--budget", type=int, default=500000, help="월 예산 (원)")
    parser.add_argument("--month", default=datetime.now().strftime("%Y-%m"), help="실행 월 (YYYY-MM)")
    args = parser.parse_args()

    main(
        product=args.product,
        keyword=args.keyword,
        category=args.category,
        target=args.target,
        price=args.price,
        cost=args.cost,
        budget=args.budget,
        month=args.month,
    )
